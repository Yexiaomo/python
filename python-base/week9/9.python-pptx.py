#!/usr/bin/env python3
# -*- coding:utf-8 -*-

## python-pptx库
'''
|接口类       |功能                     |
|-            |-                       |
|Presentation |操作ppt对象              |
|Slides       |对幻灯片进行操作          |
|Shapes       |对幻灯片的某一区域进行操作 |
|Table        |表格操作                 |
|Text         |文本操作                 |

### Presentation
|命令                    |功能              | 
|-                       |-                 |
|Presentation()          |创建一个ppt文档     |
|.slide_layouts[]        |确定幻灯片的先后顺序 |
|.slides.add_slide()     |增加一个幻灯片      |
|Slide.shape.title       |表示一个幻灯片的标题|
|Slide.shape.placeholders|表示一个幻灯片的内容|

'''
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello World!"
subtitle.text = "python-pptx is here!"

prs.save('test.pptx')