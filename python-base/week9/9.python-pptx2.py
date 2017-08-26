#!/usr/bin/env python3
# -*- coding:utf-8 -*-

## python-pptx库

## 带有图片,文本框和图形的代码

from pptx import Presentation
from pptx.util import Inches

img_path = ".\\python-pptx2.png"

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(4)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('python-pptx2.pptx')